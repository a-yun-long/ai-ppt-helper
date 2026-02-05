from openai import OpenAI  # 新加的：用来和 DeepSeek 说话
import json
import random
import re
from typing import List, Dict, Any

import streamlit as st
from pptx import Presentation  # 必须先安装：pip install python-pptx

# ==========================================
# ⚠️ 请把下面的 sk-xxxx 换成你刚才复制的那个 Key
# ==========================================
# ==========================================
# 自动判断：如果在云端就用 secrets，在本地就用你填的字符串
# ==========================================
try:
    # 尝试从云端配置读取
    DEEPSEEK_API_KEY = st.secrets["DEEPSEEK_API_KEY"]
except:
    # 如果读不到（说明在本地），就用下面这个硬编码的 Key
    DEEPSEEK_API_KEY = "sk-9e41bb1a7dbc4078bedefa87eb5aeb99" 


BASE_URL = "https://api.deepseek.com"
# ==========================================
# 1. 核心工具函数区
# ==========================================

def extract_text_from_pptx(uploaded_file):
    """从上传的PPT文件中提取所有文本"""
    try:
        prs = Presentation(uploaded_file)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    full_text.append(shape.text.strip())
        return "\n".join(full_text)
    except Exception as e:
        st.error(f"解析PPT失败: {e}")
        return ""

def normalize_text(s: str) -> str:
    s = s.strip()
    s = re.sub(r"\s+", " ", s)
    return s

def split_sentences_cn(text: str) -> List[str]:
    # 非严格分句，但够用：按 。！？；换行 切
    parts = re.split(r"[。！？；\n]+", text)
    parts = [p.strip() for p in parts if p.strip()]
    return parts

def pick_evidence(sentences: List[str], keywords: List[str], k: int = 2) -> List[str]:
    """从原文句子里选证据句：优先包含关键词的句子，否则随机挑"""
    scored = []
    for s in sentences:
        score = sum(1 for kw in keywords if kw and kw in s)
        scored.append((score, s))
    scored.sort(key=lambda x: (-x[0], -len(x[1])))
    picked = [s for score, s in scored if score > 0][:k]
    if len(picked) < k:
        pool = [s for _, s in scored if s not in picked]
        random.shuffle(pool)
        picked += pool[: (k - len(picked))]
    return picked[:k]

def simple_keywords(text: str, topn: int = 6) -> List[str]:
    # 非NLP严谨：抽取长度>=2的中文片段/英文单词，统计出现次数
    candidates = re.findall(r"[\u4e00-\u9fff]{2,}|[A-Za-z]{3,}", text)
    freq = {}
    for c in candidates:
        freq[c] = freq.get(c, 0) + 1
    items = sorted(freq.items(), key=lambda x: (-x[1], -len(x[0])))
    return [w for w, _ in items[:topn]]

# ==========================================
# 2. 模拟 AI 生成逻辑 (本地版)
# ==========================================
# ==========================================
# 2. 真正的 AI 生成逻辑 (接入 DeepSeek)
# ==========================================
def call_deepseek_generate(content: str, difficulty: str, style: str) -> Dict[str, Any]:
    """
    发送请求给 DeepSeek，让它根据 content 生成题目，并返回 JSON。
    """
    client = OpenAI(api_key=DEEPSEEK_API_KEY, base_url=BASE_URL)

    # 构造提示词 (Prompt)：告诉 AI 它是谁，我们要什么格式
    system_prompt = """
    你是一位专业的大学助教。请根据用户提供的课程内容，出具一套练习题。
    必须严格按照以下 JSON 格式返回，不要包含 Markdown 格式（如 ```json ... ```）：
    {
        "mcq": [
            {
                "id": "MCQ1",
                "stem": "题干内容...",
                "options": {"A": "...", "B": "...", "C": "...", "D": "..."},
                "answer": "A",
                "explanation": "解析...",
                "evidence": ["原文中的证据句..."]
            }
        ],
        "short": [
            {
                "id": "SA1",
                "question": "简答题问题...",
                "rubric": ["评分点1", "评分点2"],
                "evidence": ["证据句..."]
            }
        ],
        "triple": [
            {
                "id": "T1",
                "concept": {"q": "概念题...", "a": "参考答案...", "evidence": ["..."]},
                "understand": {"q": "理解题...", "a": "参考答案...", "evidence": ["..."]},
                "apply": {"q": "应用题...", "a": "参考答案...", "evidence": ["..."]}
            }
        ],
        "script_1min": {
            "title": "讲解标题",
            "sections": [{"t": "0-15s", "line": "..."}]
        }
    }
    要求：
    1. 生成 5 道单选题 (mcq)。
    2. 生成 2 道简答题 (short)。
    3. 生成 1 组三层题组 (triple)。
    4. 生成 1 分钟讲解稿 (script_1min)。
    5. 难度：{difficulty}，风格：{style}。
    6. 所有题目必须基于提供的原文，严禁瞎编。
    """

    user_prompt = f"课程内容如下：\n{content}"

    try:
        response = client.chat.completions.create(
            model="deepseek-chat",  # DeepSeek 的模型名字
            messages=[
                {"role": "system", "content": system_prompt.replace("{difficulty}", difficulty).replace("{style}", style)},
                {"role": "user", "content": user_prompt},
            ],
            temperature=0.3, # 数值越低，AI越严谨
            response_format={ 'type': 'json_object' } # 强制让它吐出 JSON
        )
        
        # 把 AI 返回的字符串变成 Python 字典
        result_json = response.choices[0].message.content
        return json.loads(result_json)

    except Exception as e:
        # 如果出错了，打印错误并在网页上报错
        print(f"DeepSeek API Error: {e}")
        # 返回一个空的结构防止程序崩溃
        return {"mcq": [], "short": [], "triple": [], "script_1min": {}, "error": str(e)}
    

def local_grade_short(rubric: List[str], user_answer: str) -> Dict[str, Any]:
    ans = normalize_text(user_answer)
    if not ans:
        return {"score": 0, "feedback": "未作答。", "missing": rubric}

    missing = []
    hit = 0
    for item in rubric:
        ks = simple_keywords(item, topn=3)
        if any(k in ans for k in ks if k):
            hit += 1
        else:
            missing.append(item)

    if hit >= max(1, len(rubric) - 0):
        score = 5
    elif hit >= max(1, len(rubric) - 1):
        score = 4
    elif hit >= max(1, len(rubric) - 2):
        score = 3
    elif hit >= 1:
        score = 2
    else:
        score = 1

    feedback = f"覆盖要点 {hit}/{len(rubric)}。建议补全：{ '；'.join(missing[:2]) }"
    return {"score": score, "feedback": feedback, "missing": missing}

def wrong_key(item: Dict[str, Any]) -> str:
    return f"{item.get('type')}::{item.get('prompt')}::{ '|'.join(item.get('evidence', [])) }"

def merge_wrong(old: List[Dict[str, Any]], new: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    m = {wrong_key(x): x for x in old}
    for x in new:
        k = wrong_key(x)
        if k not in m:
            m[k] = x
    return list(m.values())


# ==========================================
# 3. 页面 UI 主程序
# ==========================================

st.set_page_config(page_title="AI PPT学习小助手", layout="wide")

# 初始化 Session State
if "generated" not in st.session_state:
    st.session_state.generated = None
if "answers" not in st.session_state:
    st.session_state.answers = {}
if "wrongbook" not in st.session_state:
    st.session_state.wrongbook = []
if "short_grades" not in st.session_state:
    st.session_state.short_grades = {}
if "ppt_content" not in st.session_state:
    st.session_state.ppt_content = ""

st.title("AI PPT学习小助手（出题 → 纠错 → 个性化再出题）")
st.caption("离线本地版：不接外部大模型API，完整演示证据句 + 错题。")

# 侧边栏
with st.sidebar:
    st.subheader("设置")
    difficulty = st.selectbox("难度", ["基础", "中等", "偏难"], index=1)
    style = st.selectbox("风格", ["启发式", "考试导向", "更口语"], index=0)
    st.divider()
    st.write("错题本数量：", len(st.session_state.wrongbook))

    if st.button("导出错题本 JSON"):
        st.download_button(
            "点击下载",
            data=json.dumps(st.session_state.wrongbook, ensure_ascii=False, indent=2),
            file_name="wrongbook.json",
            mime="application/json",
        )

    if st.button("清空错题本"):
        st.session_state.wrongbook = []
        st.success("已清空错题本。")

# 主界面分栏
col1, col2 = st.columns([1, 1.2], gap="large")

# === 左侧栏：上传和输入 ===
with col1:
    st.subheader("① 输入课程内容")
    
    # 上传按钮
    uploaded_file = st.file_uploader("上传 PPT 课件 (自动提取文字)", type=["pptx"])
    
    # 提取逻辑
    if uploaded_file:
        # 简单判断是否是新文件，或者每次都提取
        ppt_text = extract_text_from_pptx(uploaded_file)
        if ppt_text:
            st.info(f"成功提取了 {len(ppt_text)} 个字！")
            st.session_state.ppt_content = ppt_text

    # 获取当前要显示的文字（优先用提取的 PPT 文字）
    default_text = st.session_state.get('ppt_content', "")

    # 文本框
    content = st.text_area(
        "或者直接粘贴文本", 
        height=260, 
        value=default_text,
        placeholder="粘贴文本，或者上传上方的 PPT..."
    )
    
    gen_btn = st.button("一键生成题库 + 讲解稿", type="primary", use_container_width=True)
    st.info("提示：默认每题带证据句")

    if gen_btn:
        txt = normalize_text(content)
        if len(txt) < 80:
            st.warning("文本偏短，建议至少 100 字。")
        else:
            with st.spinner("生成中…"):
                try:
                    st.session_state.generated = call_deepseek_generate(txt, difficulty, style)
                    st.session_state.answers = {}
                    st.session_state.short_grades = {}
                    st.success("生成完成！右侧可以开始做题。")
                except Exception as e:
                    st.error(f"生成失败：{e}")

# === 右侧栏：做题区域 ===
with col2:
    st.subheader("② 题库与学习闭环")
    g = st.session_state.generated

    if not g:
        st.write("先在左侧输入文本并点击生成。")
    else:
        tab_mcq, tab_sa, tab_triple, tab_script, tab_wrong = st.tabs(
            ["选择题(10)", "简答题(3)", "三层题组(3)", "1分钟讲解稿", "错题本"]
        )

        # ===== 选择题 =====
        with tab_mcq:
            for q in g.get("mcq", []):
                st.markdown(f"**{q['id']}** {q['stem']}")
                opts = q["options"]
                st.radio(
                    "选项",
                    options=["A", "B", "C", "D"],
                    format_func=lambda k: f"{k}. {opts[k]}",
                    key=f"mcq_{q['id']}",
                    horizontal=False,
                    index=None,
                )
                st.caption("证据句：")
                for ev in q.get("evidence", []):
                    st.write(f"“{ev}”")

                with st.expander("查看答案与解析（建议提交后再看）"):
                    st.write("答案：", q["answer"])
                    st.write("解析：", q["explanation"])

                st.divider()

            submit_mcq = st.button("提交作答并判分（选择题）", use_container_width=True)
            if submit_mcq:
                wrong_items = []
                for q in g.get("mcq", []):
                    user_ans = st.session_state.get(f"mcq_{q['id']}", None)
                    if user_ans != q["answer"]:
                        wrong_items.append({
                            "id": q["id"],
                            "type": "mcq",
                            "prompt": q["stem"],
                            "evidence": q.get("evidence", []),
                            "keywords": q.get("keywords", []),
                            "correct": q["answer"],
                            "user": user_ans
                        })
                st.session_state.wrongbook = merge_wrong(st.session_state.wrongbook, wrong_items)
                st.success(f"判分完成：新增错题 {len(wrong_items)} 道；错题本共 {len(st.session_state.wrongbook)} 道。")

        # ===== 简答题 =====
        with tab_sa:
            grade_btn = st.button("提交作答并评分（简答题）", type="primary", use_container_width=True)

            for q in g.get("short", []):
                st.markdown(f"**{q['id']}** {q['question']}")
                st.text_area("你的答案", key=f"sa_{q['id']}", height=120)
                st.caption("Rubric：")
                for r in q.get("rubric", []):
                    st.write(f"- {r}")
                st.caption("证据句：")
                for ev in q.get("evidence", []):
                    st.write(f"“{ev}”")

                if q["id"] in st.session_state.short_grades:
                    gr = st.session_state.short_grades[q["id"]]
                    st.write(f"得分：**{gr['score']} / 5**")
                    st.write("反馈：", gr.get("feedback", ""))
                    if gr.get("missing"):
                        st.write("缺失：", "；".join(gr["missing"]))

                st.divider()

            if grade_btn:
                wrong_items = []
                for q in g.get("short", []):
                    user_answer = normalize_text(st.session_state.get(f"sa_{q['id']}", ""))
                    gr = local_grade_short(q["rubric"], user_answer)

                    st.session_state.short_grades[q["id"]] = gr

                    if gr["score"] <= 2:
                        wrong_items.append({
                            "id": q["id"],
                            "type": "short",
                            "prompt": q["question"],
                            "evidence": q.get("evidence", []),
                            "keywords": [],
                            "correct": q.get("rubric", []),
                            "user": user_answer,
                            "grade": gr
                        })

                st.session_state.wrongbook = merge_wrong(st.session_state.wrongbook, wrong_items)
                st.success(f"评分完成：新增错题 {len(wrong_items)} 道；错题本共 {len(st.session_state.wrongbook)} 道。")
                st.rerun()

        # ===== 三层题组 =====
        with tab_triple:
            for t in g.get("triple", []):
                st.markdown(f"### {t['id']}")
                for name, key in [("概念题（记忆）", "concept"), ("理解题（解释）", "understand"), ("应用题（小案例）", "apply")]:
                    blk = t[key]
                    st.markdown(f"**{name}**：{blk['q']}")
                    st.write("参考：", blk["a"])
                    st.caption("证据句：")
                    for ev in blk.get("evidence", []):
                        st.write(f"“{ev}”")
                st.divider()

        #讲解稿
        with tab_script:
            s = g.get("script_1min", {})
            st.markdown(f"### {s.get('title', '1分钟讲解稿')}")
            for sec in s.get("sections", []):
                st.markdown(f"**{sec['t']}** {sec['line']}")
            st.download_button(
                "下载讲解稿（txt）",
                data="\n".join([f"{x['t']} {x['line']}" for x in s.get("sections", [])]),
                file_name="script_1min.txt"
            )

        #错题本
        with tab_wrong:
            wb = st.session_state.wrongbook
            if not wb:
                st.write("还没有错题。做题并提交后会自动加入。")
            else:
                for i, item in enumerate(wb):
                    st.markdown(f"**{i+1}. {item['id']}** ({item['type']})")
                    st.write(item["prompt"])
                    st.caption(f"你的答案：{item.get('user')}")
                    st.caption(f"正确参考：{item.get('correct')}")
                    st.caption("证据句：")
                    for ev in item.get("evidence", []):
                        st.write(f"“{ev}”")
                    if st.button("删除这题", key=f"del_{i}"):
                        wb2 = wb[:i] + wb[i+1:]
                        st.session_state.wrongbook = wb2
                        st.rerun()
                    st.divider()

            st.markdown("### 一键针对错题再出题（变式训练）")
            if st.button("生成变式题（覆盖选择+简答）", type="primary", use_container_width=True):
                if not content.strip():
                    st.warning("需要原文内容才能生成变式题。")
                elif not st.session_state.wrongbook:
                    st.warning("错题本为空。")
                else:
                    with st.spinner("生成变式题中…"):
                        txt = normalize_text(content)
                        gen = local_generate(txt, difficulty, style)

                        # 简单个性化：根据错题关键词强化
                        allk = []
                        for it in st.session_state.wrongbook[-8:]:
                            allk += it.get("keywords", [])
                        if allk:
                            gen["meta"]["topic"] = "薄弱点训练：" + max(set(allk), key=allk.count)

                        # 只保留少量，强调再训练
                        gen["mcq"] = gen["mcq"][:3]
                        gen["short"] = gen["short"][:3]
                        st.session_state.generated = gen

                        st.session_state.answers = {}
                        st.session_state.short_grades = {}
                        st.success("变式题已生成：选择题3道 + 简答题3道。去左边 tabs 做题。")
                        st.rerun()